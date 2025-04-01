from smolagents import tool


@tool
def read_pdf_file(file_path: str) -> str:
    """
    This function reads the first three pages of a PDF file and returns its content as a string.
    Args:
        file_path: The path to the PDF file.
    Returns:
        A string containing the content of the PDF file.
    """
    from pypdf import PdfReader

    content = ""
    reader = PdfReader(file_path)
    print(len(reader.pages))
    pages = reader.pages[:3]
    for page in pages:
        content += page.extract_text()
    return content


@tool
def read_docx_file(file_path: str) -> str:
    """
    This function reads the content of a DOCX file and returns it as a string.
    Args:
        file_path: The path to the DOCX file.
    Returns:
        A string containing the content of the DOCX file.
    """
    from docx import Document

    content = ""
    document = Document(file_path)
    for paragraph in document.paragraphs:
        content += paragraph.text + "\n"
    return content


@tool
def read_pptx_file(file_path: str) -> str:
    """
    This function reads the content of a PPTX file and returns it as a string.
    Args:
        file_path: The path to the PPTX file.
    Returns:
        A string containing the content of the PPTX file.
    """
    from pptx import Presentation

    content = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + "\n"
    return content


@tool
def read_text_file(file_path: str) -> str:
    """
    This function reads the content of a plain text file and returns it as a string.
    Args:
        file_path: The path to the text file.
    Returns:
        A string containing the content of the text file.
    """
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()
    return content
